import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from io import BytesIO

# Use Streamlit's Secret Manager to stay secure
api_key = st.secrets["GEMINI_API_KEY"]
genai.configure(api_key=api_key)
model = genai.GenerativeModel('gemini-1.5-flash')

st.title("Kobby's AI Slide Generator")

topic = st.text_input("What is your presentation about?")

if st.button("Generate Slides"):
    with st.spinner("AI is writing..."):
        # AI writes the content
        response = model.generate_content(f"Create a 3-slide outline for: {topic}")
        
        # Build the PPTX
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = topic
        slide.placeholders[1].text = response.text
        
        # Save to memory
        buf = BytesIO()
        prs.save(buf)
        
        st.success("Slides Ready!")
        st.download_button("📥 Download PowerPoint", buf.getvalue(), "presentation.pptx")
